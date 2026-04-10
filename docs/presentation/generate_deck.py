"""Generate the CryptoFlow EDPO presentation from the outline, using the HSG template.

Run:
    python3 -m venv /tmp/pptx-venv && /tmp/pptx-venv/bin/pip install python-pptx
    /tmp/pptx-venv/bin/python docs/presentation/generate_deck.py

Output: docs/presentation/cryptoflow.pptx
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.oxml.ns import qn
from PIL import Image


REPO = Path(__file__).resolve().parents[2]
TEMPLATE = REPO / "docs/presentation/HSG-Template.pptx"
OUTPUT = REPO / "docs/presentation/cryptoflow.pptx"
FIG_REPORT = REPO / "docs/report/gallus-hsg/content/chapters/figures"
FIG_LOCAL = REPO / "docs/presentation/figures"

CONTEXT_MAP = FIG_REPORT / "context-map.jpeg"
DEPLOYMENT_OVERVIEW = FIG_LOCAL / "deployment-overview.png"
ONBOARDING_BPMN = FIG_REPORT / "userOnboarding-bpmn.png"
PLACE_ORDER_BPMN = FIG_REPORT / "placeOrder-bpmn.png"


# --- Layout indices (from template inspection) ---
LAYOUT_TITLE = 0            # Titel Slide 1
LAYOUT_CONTENT = 7          # Titel und Content  (title + single body)
LAYOUT_TWO_CONTENT = 6      # Titel zwei Content (title + two side-by-side bodies)
LAYOUT_TITLE_ONLY = 10      # Nur Titel


def delete_all_slides(prs: Presentation) -> None:
    """Remove every slide reference from the template while keeping masters/layouts."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def get_placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_title(slide, text: str) -> None:
    ph = slide.shapes.title
    if ph is None:
        ph = get_placeholder(slide, 0)
    if ph is None:
        raise RuntimeError("no title placeholder on slide")
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def set_bullets(placeholder, items, *, base_size: int = 18) -> None:
    """Fill a placeholder's text frame with bullets.

    `items` is a list of either strings (level 0) or (text, level) tuples.
    """
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.text = text
        # Size shrinks one step per nesting level. Template handles bullet glyphs.
        size = base_size - 2 * level
        for run in p.runs:
            run.font.size = Pt(size)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def add_picture_in_placeholder(slide, placeholder, image_path: Path) -> None:
    """Replace the given placeholder with a picture that fits its bbox, preserving aspect ratio."""
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    # Remove the content placeholder shape so it doesn't overlap the picture.
    sp = placeholder._element
    sp.getparent().remove(sp)
    # Compute fit-preserving box.
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = width / height
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # image is wider — fit to width
        new_w = width
        new_h = int(width / img_ratio)
    else:
        new_h = height
        new_w = int(height * img_ratio)
    new_left = left + (width - new_w) // 2
    new_top = top + (height - new_h) // 2
    slide.shapes.add_picture(str(image_path), new_left, new_top, new_w, new_h)


def add_slide(prs, layout_index: int):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


# ---------------- Slide builders ----------------

def slide_1_title(prs):
    s = add_slide(prs, LAYOUT_TITLE)
    set_title(s, "CryptoFlow\nEvent-Driven & Process-Oriented Crypto Portfolio Platform")
    body = get_placeholder(s, 10)
    if body is not None:
        tf = body.text_frame
        tf.clear()
        tf.paragraphs[0].text = "EDPO Project FS26"
        p = tf.add_paragraph()
        p.text = "St.Gallen  |  Team of 2"
    set_notes(
        s,
        "Open confidently. We built CryptoFlow to exercise the full set of EDPO lecture "
        "concepts end to end: two communication channels, five microservices, two saga "
        "patterns. This slide is just the name card — keep it to ten seconds.",
    )


def slide_2_what_is(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "What is CryptoFlow?")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Crypto portfolio simulation platform — not a production exchange",
            "Users register, observe live market prices, manage portfolios, place simulated trades",
            "Case study for applying EDPO concepts in one coherent end-to-end system",
        ],
        base_size=20,
    )
    set_notes(
        s,
        "Keep this brief. The point is that we chose a domain rich enough to require "
        "multiple bounded contexts, long-running processes, and eventual consistency — "
        "all things we wanted to demonstrate.",
    )


def slide_3_contexts(prs):
    s = add_slide(prs, LAYOUT_TWO_CONTENT)
    set_title(s, "Domain — Five Bounded Contexts")
    left = get_placeholder(s, 1)
    set_bullets(
        left,
        [
            "Market Data — upstream price feed",
            "Portfolio — holdings & valuation",
            "Trading — order lifecycle",
            "User Identity — accounts & confirmation",
            "Onboarding — saga orchestrator",
            ("Upstream/downstream relationships", 1),
            ("Partnership: User ↔ Portfolio (compensation)", 1),
            ("Shared kernel: shared-events module", 1),
        ],
        base_size=16,
    )
    right = get_placeholder(s, 2)
    add_picture_in_placeholder(s, right, CONTEXT_MAP)
    set_notes(
        s,
        "Emphasize that the bounded contexts drove the service decomposition. Each "
        "context is owned by exactly one service. Cross-context communication is always "
        "asynchronous.",
    )


def slide_4_architecture(prs):
    s = add_slide(prs, LAYOUT_TWO_CONTENT)
    set_title(s, "Architecture Overview")
    left = get_placeholder(s, 1)
    set_bullets(
        left,
        [
            "5 Spring Boot microservices in Docker Compose",
            "Two communication channels: Kafka + Zeebe",
            "No synchronous inter-service calls for domain ops",
            "REST only for external clients",
            "Database-per-service (3× PostgreSQL)",
            ("Market Data and Onboarding are stateless", 1),
        ],
        base_size=16,
    )
    right = get_placeholder(s, 2)
    add_picture_in_placeholder(s, right, DEPLOYMENT_OVERVIEW)
    set_notes(
        s,
        "This is the 30-second mental model. Two channels, five services, no shared "
        "databases. Everything else in the presentation builds on this.",
    )


def slide_5_event_driven(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Event-Driven Communication (High Level)")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Kafka as the sole inter-service event bus (ADR-0001)",
            "6 topics: price ticks, approved orders, user confirmations, compensation (×2), DLT",
            "Event-Carried State Transfer in 4 places",
            ("Price cache, portfolio updates, compensation events, replicated user validation", 1),
            "Key config: acks=all, auto.offset.reset=earliest, manual commit, DLT after 3 attempts",
        ],
        base_size=18,
    )
    set_notes(
        s,
        "We won't deep-dive into event-driven architecture here — that's the focus of the "
        "next presentation. The key takeaway: services communicate through published facts, "
        "not remote commands. ECST lets each service maintain its own local view without "
        "querying others.",
    )


def slide_6_kafka_experiments(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Kafka Experiments — Assignment 1 (Summary)")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Four experiments on multi-broker clusters validating Kafka reliability",
            "acks=0: eventID=76 permanently lost — producer cannot detect the loss",
            "acks=1: eventID=71 ACKed by leader but lost when leader crashed before replication",
            "acks=all: ~8.9s availability gap during failover, zero data loss",
            "Offset reset: earliest replays retained history; latest silently skips it",
            "Findings directly informed our production config choices (ADR-0001)",
        ],
        base_size=16,
    )
    set_notes(
        s,
        "Show the producer/consumer log snippets briefly. The main insight: durability is "
        "not a Kafka default — it requires explicit configuration. Our single-broker dev "
        "setup means acks=all is equivalent to acks=1 locally; we accept this as a conscious "
        "dev-environment trade-off.",
    )


def slide_7_why_camunda(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Why Camunda 8 / Zeebe?")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Transition: from events to processes",
            "Two workflows need durable waiting states — onboarding and order placement",
            "Why not Camunda 7 / Operaton?",
            ("Embedded engine stores process state in the application DB", 1),
            ("Many concurrent placeOrder waits would need a shared polling table", 1),
            ("Kafka and SMTP require custom application-level client code", 1),
            "Camunda 8 / Zeebe advantages",
            ("Partitioned log — each instance waits independently, no job-lock contention", 1),
            ("Connector templates — Kafka + SMTP inside BPMN, zero notification code", 1),
            ("Camunda Operate — runtime visibility without custom dashboards", 1),
        ],
        base_size=14,
    )
    set_notes(
        s,
        "This is not a generic 'Camunda 8 is better' argument. These are project-specific "
        "reasons. The partitioned log matters because placeOrder has a non-deterministic "
        "wait — many orders can be open simultaneously. The connector templates matter "
        "because our flows include email notifications.",
    )


def slide_8_two_sagas(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Two Saga Patterns — Driven by Wait Semantics")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Two orchestrated saga styles — choice driven by nature of the wait",
            "Onboarding → Parallel Saga (aeo)",
            ("Deterministic wait (user clicks confirmation link)", 1),
            ("Spans two bounded contexts — User + Portfolio", 1),
            ("Needs all-or-nothing consistency", 1),
            "Order Placement → Fairy Tale Saga (seo)",
            ("Non-deterministic wait (price match from market feed)", 1),
            ("Stays inside one bounded context — Trading", 1),
            ("Accepts eventual consistency with Portfolio", 1),
        ],
        base_size=14,
    )
    set_notes(
        s,
        "This was one of the most instructive decisions. The wait semantics dictated the "
        "coupling model, which dictated the saga style. A deterministic wait with "
        "cross-context scope → Parallel. A non-deterministic wait within one context → "
        "Fairy Tale.",
    )


def slide_9_onboarding_process(prs):
    s = add_slide(prs, LAYOUT_TWO_CONTENT)
    set_title(s, "Onboarding Process — Parallel Saga")
    left = get_placeholder(s, 1)
    set_bullets(
        left,
        [
            "prepareUserWorker: generates userId, PENDING link",
            "Camunda email connector sends confirmation",
            "Event-based gateway: click vs. 1-min timeout",
            "Parallel fan-out: userCreationWorker + portfolioCreationWorker",
            "Both must succeed → process completes",
        ],
        base_size=14,
    )
    right = get_placeholder(s, 2)
    add_picture_in_placeholder(s, right, ONBOARDING_BPMN)
    set_notes(
        s,
        "Point to the BPMN diagram as you walk through. Emphasize: three services "
        "participate but only one (onboarding-service) owns the BPMN. The workers are "
        "stateless — they execute local transactions and report back.",
    )


def slide_10_compensation(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Onboarding — Compensation & Flag-Driven Completion")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Challenge: one creation succeeds, the other fails — what now?",
            "Flag-driven completion (ADR-0012)",
            ("Workers always complete Zeebe jobs, report outcome via flags", 1),
            ("isUserCreated / isPortfolioCreated drive gateway branches", 1),
            ("BPMN errors would short-circuit and bypass compensation gateways", 1),
            "Compensation (ADR-0011)",
            ("Exclusive gateways inspect flags → compensation handlers", 1),
            ("Delete local entity + publish compensation request for peer", 1),
            ("User/Portfolio CompensationRequestedEvent carries full payload", 1),
            ("Idempotent deletes — safe under at-least-once delivery", 1),
        ],
        base_size=14,
    )
    set_notes(
        s,
        "This is where theory meets implementation. The flag-driven pattern is subtle but "
        "essential — without it, the BPMN gateway branches become unreachable. The "
        "compensation events use ECST so no synchronous callback is needed.",
    )


def slide_11_dedicated_service(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Dedicated Onboarding Service (ADR-0010)")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Originally the BPMN lived inside user-service",
            "Problem: orchestrator coordinates work across two bounded contexts",
            ("Coupling + mixed responsibilities violates DDD boundaries", 1),
            "Solution: extract into a dedicated onboarding-service",
            ("Owns the BPMN process and Zeebe deployment", 1),
            ("No persistent business data — workflow state lives in Zeebe", 1),
            ("User and portfolio services keep their domain autonomy", 1),
            "Trade-off: 5 services instead of 4, but clean separation",
        ],
        base_size=16,
    )
    set_notes(
        s,
        "This was an ADR-driven decision mid-project. The original design violated bounded "
        "context boundaries. The dedicated service is the most important structural "
        "consequence of applying the saga pattern correctly.",
    )


def slide_12_place_order(prs):
    s = add_slide(prs, LAYOUT_TWO_CONTENT)
    set_title(s, "Place Order Process — Fairy Tale Saga")
    left = get_placeholder(s, 1)
    set_bullets(
        left,
        [
            "order-processing-worker validates user locally, creates PENDING tx",
            "Event-based gateway: price match vs. 1-min timeout",
            "approveOrderWorker writes APPROVED + outbox row in one transaction",
            "publishOrderApprovedWorker publishes to Kafka → email",
            "Portfolio update happens independently via Kafka",
        ],
        base_size=14,
    )
    right = get_placeholder(s, 2)
    add_picture_in_placeholder(s, right, PLACE_ORDER_BPMN)
    set_notes(
        s,
        "Key distinction from onboarding: the approval is terminal in the trading context. "
        "Portfolio propagation is a downstream consequence, not a saga step. This is "
        "exactly the Fairy Tale Saga trade-off: synchronous within the context, eventual "
        "across contexts.",
    )


def slide_13_reliability(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Reliability Patterns in the Trading Flow")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Three patterns protecting the order execution path",
            "Transactional Outbox (ADR-0014)",
            ("Approval + outbox row in one DB transaction", 1),
            ("Scheduler republishes stale rows after crashes", 1),
            "Idempotent Consumer (ADR-0016)",
            ("processed_transaction UNIQUE constraint on transactionId", 1),
            ("Found during integration: 2 BTC order → 4 BTC in portfolio", 1),
            "Human Escalation (ADR-0018)",
            ("Deterministic failures throw typed BPMN errors", 1),
            ("Boundary events route to ops user task in Tasklist", 1),
        ],
        base_size=14,
    )
    set_notes(
        s,
        "These three patterns form a chain: the outbox guarantees the event is published, "
        "the idempotent consumer guarantees it's applied exactly once, and human escalation "
        "handles the cases that automated retries cannot fix. The idempotency bug was "
        "discovered accidentally during integration — a real validation of the pattern's "
        "necessity.",
    )


def slide_14_read_model(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Replicated Read-Model for Autonomous Validation")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Problem: reject orders from unconfirmed users without calling user-service",
            "Solution (ADR-0017): transaction-service keeps its own table of confirmed users",
            ("user-service publishes UserConfirmedEvent to a log-compacted topic", 1),
            ("transaction-service consumes and upserts into local DB", 1),
            ("Order validation = local read, sub-millisecond, no network call", 1),
            "Trade-off: brief eventual consistency window",
            ("Acceptable — a newly confirmed user won't place an order immediately", 1),
        ],
        base_size=16,
    )
    set_notes(
        s,
        "This is a focused use of CQRS: the write model stays in user-service, the trading "
        "context keeps only the projection it needs. Log compaction ensures the read-model "
        "survives cold starts.",
    )


def slide_15_challenges(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Challenges & Lessons Learned")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "BPMN learning curve — demanding upfront, reduced implementation effort after",
            "Concepts intuitive in theory, complex in practice",
            ("Orchestration, compensation, ECST multiply downstream decisions", 1),
            "Saga selection only became clear by analyzing the nature of the wait",
            "Integration issues surface late — deserialization, correlation keys, payloads",
            "Two people, five services — ambitious scope, ADR-0010 added overhead",
        ],
        base_size=16,
    )
    set_notes(
        s,
        "Be honest about the challenges. The key message: the patterns work, but applying "
        "them requires many more decisions than the theory suggests. ADRs were essential "
        "for keeping those decisions traceable.",
    )


def slide_16_demo(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Live Demo")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "Demo 1 — Happy-path onboarding",
            ("Tasklist form → confirmation email → click → user + portfolio created", 1),
            ("Camunda Operate shows the completed process instance", 1),
            "Demo 2 — Happy-path order placement",
            ("Tasklist order → price match → APPROVED → portfolio updated", 1),
            ("Notification email received", 1),
            "Optional if time permits: timeout scenario, Kafka UI, Operate history",
        ],
        base_size=14,
    )
    set_notes(
        s,
        "Keep the demo tight — 2-3 minutes max. Have the Docker Compose stack running "
        "before the presentation starts. If something breaks live, switch to Operate to "
        "show the instance state — that's still a valid demo of observability.",
    )


def slide_17_summary(prs):
    s = add_slide(prs, LAYOUT_CONTENT)
    set_title(s, "Summary & Key Takeaways")
    body = get_placeholder(s, 1)
    set_bullets(
        body,
        [
            "5 microservices, 2 channels (Kafka + Zeebe), 11 EDPO concepts implemented",
            "Two saga patterns chosen by wait semantics, not by preference",
            "Flag-driven completion keeps compensation paths reachable",
            "Transactional outbox + idempotent consumer + human escalation as a chain",
            "Dedicated orchestration service preserves bounded context autonomy",
            "20 ADRs documenting every architectural decision with rationale",
        ],
        base_size=16,
    )
    set_notes(
        s,
        "This is the wrap-up. Reiterate the most distinctive architectural choices. Then "
        "open for Q&A.",
    )


def slide_18_qa(prs):
    s = add_slide(prs, LAYOUT_TITLE_ONLY)
    set_title(s, "Questions?")
    # Add a small subtitle-ish text box under the title with backup hints.
    left = Inches(0.4)
    top = Inches(2.2)
    width = Inches(9.2)
    height = Inches(2.8)
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = "Anticipated backup topics"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(18)
        run.font.bold = True
    for txt in (
        "Why orchestration rather than choreography?",
        "Why JSON over Avro for Kafka serialization?",
        "How would you scale this to production?",
        "What would you do differently next time?",
    ):
        p = tf.add_paragraph()
        p.text = "•  " + txt
        for run in p.runs:
            run.font.size = Pt(16)
    set_notes(
        s,
        "Open the floor. If the audience is quiet, prime the pump with the first backup "
        "topic yourself.",
    )


# ---------------- Main ----------------

def main() -> None:
    assert TEMPLATE.exists(), f"missing template: {TEMPLATE}"
    for p in (CONTEXT_MAP, DEPLOYMENT_OVERVIEW, ONBOARDING_BPMN, PLACE_ORDER_BPMN):
        assert p.exists(), f"missing figure: {p}"

    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    builders = [
        slide_1_title,
        slide_2_what_is,
        slide_3_contexts,
        slide_4_architecture,
        slide_5_event_driven,
        slide_6_kafka_experiments,
        slide_7_why_camunda,
        slide_8_two_sagas,
        slide_9_onboarding_process,
        slide_10_compensation,
        slide_11_dedicated_service,
        slide_12_place_order,
        slide_13_reliability,
        slide_14_read_model,
        slide_15_challenges,
        slide_16_demo,
        slide_17_summary,
        slide_18_qa,
    ]
    for fn in builders:
        fn(prs)

    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT}  ({OUTPUT.stat().st_size // 1024} KB, {len(builders)} slides)")


if __name__ == "__main__":
    main()
